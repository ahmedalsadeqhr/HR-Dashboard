import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 51Talk Brand Refresh — March 2026
NAVY       = RGBColor(0x1B, 0x2A, 0x5C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GOLD       = RGBColor(0xF5, 0xC4, 0x30)
LIGHT_GRAY = RGBColor(0xB0, 0xBE, 0xC5)
RED        = RGBColor(0xE5, 0x39, 0x35)
BLUE       = RGBColor(0x3D, 0x8E, 0xF0)
AMBER      = RGBColor(0xF5, 0xC4, 0x30)
GREEN      = RGBColor(0x43, 0xA0, 0x47)
DARK_CARD  = RGBColor(0x16, 0x20, 0x44)

W = Inches(13.33)
H = Inches(7.5)

MSO_RECT    = 1
MSO_OVAL    = 9
MSO_DIAMOND = 4


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height,
        bg_color=None, text='', font_size=18,
        bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        wrap=True, line_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    else:
        txBox.fill.background()
    if line_color:
        txBox.line.color.rgb = line_color
        txBox.line.width = Pt(1)
    else:
        txBox.line.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = 'Calibri'
    return txBox


def shape_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1.5):
    s = slide.shapes.add_shape(MSO_RECT, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    return s


def accent_line(slide, top, color=GOLD, left=Inches(0.5), width=Inches(12.33)):
    ln = slide.shapes.add_shape(MSO_RECT, left, top, width, Pt(3))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    return ln


def star_deco(slide, left, top, size, color=GOLD, alpha=180):
    """4-pointed sparkle using overlapping thin rectangles."""
    cx = left + size / 2
    cy = top  + size / 2
    thin = size * 0.18
    for w, h in [(size, thin), (thin, size)]:
        r = slide.shapes.add_shape(MSO_RECT,
            cx - w/2, cy - h/2, w, h)
        r.fill.solid()
        r.fill.fore_color.rgb = color
        r.line.fill.background()


# ---------------------------------------------------------------------------
# Slide 1 — Cover
# ---------------------------------------------------------------------------
def slide_01_cover(prs):
    s = blank_slide(prs)
    bg(s)

    # Decorative stars
    for (lx, ly, sz) in [
        (Inches(10.8), Inches(0.4), Inches(0.35)),
        (Inches(12.0), Inches(1.1), Inches(0.22)),
        (Inches(11.5), Inches(2.0), Inches(0.18)),
        (Inches(0.3),  Inches(5.8), Inches(0.20)),
        (Inches(1.5),  Inches(6.8), Inches(0.15)),
        (Inches(12.5), Inches(5.5), Inches(0.25)),
    ]:
        star_deco(s, lx, ly, sz)

    # Left gold accent bar
    shape_rect(s, Inches(0), Inches(0), Inches(0.1), H, GOLD)

    # Yellow circle background decoration (top right)
    circ = s.shapes.add_shape(MSO_OVAL, Inches(10.2), Inches(-1.0), Inches(4.5), Inches(4.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = DARK_CARD
    circ.line.color.rgb = GOLD; circ.line.width = Pt(2)

    # 51TALK label
    box(s, Inches(0.5), Inches(1.4), Inches(4), Inches(0.45),
        text='51TALK EGYPT', font_size=13, bold=True, color=GOLD)

    # Main title
    box(s, Inches(0.5), Inches(1.9), Inches(10), Inches(1.3),
        text='Work Review Report', font_size=48, bold=True, color=WHITE)

    box(s, Inches(0.5), Inches(3.15), Inches(9), Inches(0.65),
        text='Probation Period Review', font_size=24, color=GOLD)

    accent_line(s, Inches(3.95), width=Inches(6))

    # Details
    box(s, Inches(0.5), Inches(4.25), Inches(5), Inches(0.5),
        text='Ahmed Al-Sadeq', font_size=20, bold=True, color=WHITE)
    box(s, Inches(0.5), Inches(4.75), Inches(7), Inches(0.45),
        text='HR Business Partner  |  51Talk Egypt', font_size=15, color=LIGHT_GRAY)
    box(s, Inches(0.5), Inches(5.2), Inches(7), Inches(0.45),
        text='December 9, 2025  –  March 9, 2026', font_size=14, color=LIGHT_GRAY)

    return s


# ---------------------------------------------------------------------------
# Slide 2 — Self Introduction
# ---------------------------------------------------------------------------
def slide_02_intro(prs):
    s = blank_slide(prs)
    bg(s)

    box(s, Inches(0.5), Inches(0.22), Inches(4), Inches(0.38),
        text='01  /  Self Introduction', font_size=11, color=GOLD, bold=True)
    box(s, Inches(0.5), Inches(0.6), Inches(12), Inches(0.72),
        text='Who I Am', font_size=34, bold=True, color=WHITE)
    accent_line(s, Inches(1.42))

    # Left — key fact cards
    cards = [
        ('Date of Joining',      'December 9, 2025'),
        ('Years of Experience',  '10 Years'),
        ('Current Role',         'HR Business Partner'),
        ('Department',           'Human Resources — 51Talk Egypt'),
    ]
    y = Inches(1.62)
    for label, value in cards:
        shape_rect(s, Inches(0.5), y, Inches(4.7), Inches(0.75), DARK_CARD)
        box(s, Inches(0.65), y + Pt(5), Inches(4.4), Inches(0.28),
            text=label, font_size=10, color=GOLD, bold=True)
        box(s, Inches(0.65), y + Pt(24), Inches(4.4), Inches(0.35),
            text=value, font_size=14, color=WHITE)
        y += Inches(0.85)

    # Right — summary
    box(s, Inches(5.6), Inches(1.62), Inches(7.2), Inches(0.38),
        text='Professional Summary', font_size=14, bold=True, color=GOLD)

    box(s, Inches(5.6), Inches(2.05), Inches(7.2), Inches(1.55),
        text=(
            'HR professional with 10 years of experience spanning all HR functions — '
            'from talent acquisition, onboarding, and employee relations to HR operations, '
            'legal compliance, and HR technology. Proven track record of managing '
            'high-volume HR operations, building systems that scale, and driving '
            'people-focused initiatives that create measurable business impact.'
        ),
        font_size=13, color=LIGHT_GRAY, wrap=True)

    box(s, Inches(5.6), Inches(3.72), Inches(5), Inches(0.38),
        text='Key Responsibilities', font_size=14, bold=True, color=GOLD)

    resp = [
        '▸  Employee lifecycle management (onboarding to offboarding)',
        '▸  Employee relations and case management',
        '▸  HR systems implementation and automation',
        '▸  Legal compliance (social insurance, labor law)',
        '▸  HR data, reporting, and dashboards',
    ]
    y2 = Inches(4.15)
    for r in resp:
        box(s, Inches(5.6), y2, Inches(7.2), Inches(0.36),
            text=r, font_size=13, color=WHITE)
        y2 += Inches(0.37)

    return s


# ---------------------------------------------------------------------------
# Slide 3 — Performance Summary
# ---------------------------------------------------------------------------
def slide_03_performance(prs):
    s = blank_slide(prs)
    bg(s)

    box(s, Inches(0.5), Inches(0.22), Inches(5), Inches(0.38),
        text='02  /  Performance Summary', font_size=11, color=GOLD, bold=True)
    box(s, Inches(0.5), Inches(0.6), Inches(12), Inches(0.72),
        text='Impact at a Glance', font_size=34, bold=True, color=WHITE)
    accent_line(s, Inches(1.42))

    # Top row — 3 stat cards
    stats = [
        ('100', 'New Hires Onboarded',
         'End-to-end onboarding for 100\nemployees within 3 months'),
        ('50',  'ER Cases Handled',
         'Managed and resolved 50\nemployee relations cases'),
        ('3',   'Legal & Compliance Cases',
         '1 social insurance investigation\n2 re-hire violation cases resolved'),
    ]
    x = Inches(0.5)
    for num, title, desc in stats:
        shape_rect(s, x, Inches(1.62), Inches(4.0), Inches(2.1), DARK_CARD,
                   line_color=GOLD, line_width=1)
        box(s, x + Inches(0.2), Inches(1.72), Inches(3.6), Inches(0.95),
            text=num, font_size=54, bold=True, color=GOLD)
        box(s, x + Inches(0.2), Inches(2.6), Inches(3.6), Inches(0.35),
            text=title, font_size=13, bold=True, color=WHITE)
        box(s, x + Inches(0.2), Inches(2.96), Inches(3.6), Inches(0.55),
            text=desc, font_size=11, color=LIGHT_GRAY, wrap=True)
        x += Inches(4.25)

    # Bottom row — 2 initiative cards
    initiatives = [
        ('3', 'HR Systems Built & Automated',
         '• Leave Email Automation (Lark integration)\n'
         '• Attendance & HC Dashboard\n'
         '• Live Quiz System (600+ concurrent players)'),
        ('1', 'Major Initiative In Progress',
         'iTalent ESS Portal — transitioning leave management\n'
         'from email to fully digital self-service system'),
    ]
    x = Inches(0.5)
    for num, title, desc in initiatives:
        shape_rect(s, x, Inches(3.95), Inches(6.15), Inches(3.15), DARK_CARD,
                   line_color=GOLD, line_width=1)
        box(s, x + Inches(0.2), Inches(4.05), Inches(5.7), Inches(0.85),
            text=num, font_size=46, bold=True, color=GOLD)
        box(s, x + Inches(0.2), Inches(4.82), Inches(5.7), Inches(0.38),
            text=title, font_size=14, bold=True, color=WHITE)
        box(s, x + Inches(0.2), Inches(5.22), Inches(5.7), Inches(1.6),
            text=desc, font_size=12, color=LIGHT_GRAY, wrap=True)
        x += Inches(6.55)

    return s


# ---------------------------------------------------------------------------
# Slide 4 — Value Practice (STAR)
# ---------------------------------------------------------------------------
def slide_04_values(prs):
    s = blank_slide(prs)
    bg(s)

    box(s, Inches(0.5), Inches(0.22), Inches(5), Inches(0.38),
        text='03  /  Value Practice', font_size=11, color=GOLD, bold=True)
    box(s, Inches(0.5), Inches(0.6), Inches(12.5), Inches(0.72),
        text='Living the 51Talk Values — STAR Method', font_size=30, bold=True, color=WHITE)
    accent_line(s, Inches(1.42))

    values = [
        (RED,   'Game Changer',
         'Leave management was fully email-based — slow and error-prone.',
         'Find and implement digital HR solutions proactively.',
         'Built Leave Automation, Attendance Dashboard & iTalent implementation plan.',
         'Reduced manual effort; laid foundation for a fully digital HR operation.'),
        (BLUE,  'Customer Focus',
         '51Talk Egypt needed structured onboarding to support rapid growth.',
         'Ensure every new hire has a smooth, consistent experience.',
         'Managed end-to-end onboarding for 100 employees in 3 months.',
         '100 employees fully integrated — zero onboarding gaps reported.'),
        (GREEN, 'Passion',
         'ER cases demand time, sensitivity, and deep HR expertise.',
         'Handle all cases professionally while maintaining employee trust.',
         'Managed 50 ER cases with full documentation and fair investigation.',
         'All cases resolved within policy — workplace stability maintained.'),
        (AMBER, 'Teamwork',
         'Compliance cases required coordinated cross-functional effort.',
         'Navigate legal matters in collaboration with legal & finance teams.',
         'Led 3 investigations, coordinating with legal and finance departments.',
         'All cases resolved in full compliance with Egyptian labor law.'),
    ]

    positions = [
        (Inches(0.4),  Inches(1.6)),
        (Inches(6.75), Inches(1.6)),
        (Inches(0.4),  Inches(4.45)),
        (Inches(6.75), Inches(4.45)),
    ]

    star_labels = ['S', 'T', 'A', 'R']
    result_colors = [LIGHT_GRAY, LIGHT_GRAY, LIGHT_GRAY, GOLD]

    for (color, value_name, s_txt, t_txt, a_txt, r_txt), (cx, cy) in zip(values, positions):
        # Card
        shape_rect(s, cx, cy, Inches(6.1), Inches(2.7), DARK_CARD,
                   line_color=color, line_width=1.5)
        # Header
        shape_rect(s, cx, cy, Inches(6.1), Inches(0.42), color)
        box(s, cx + Inches(0.15), cy + Pt(5), Inches(5.8), Inches(0.32),
            text=value_name, font_size=14, bold=True, color=WHITE)

        # STAR rows
        star_data  = [s_txt, t_txt, a_txt, r_txt]
        for i, (lbl, data) in enumerate(zip(star_labels, star_data)):
            row_y = cy + Inches(0.5) + i * Inches(0.54)
            # Badge
            shape_rect(s, cx + Inches(0.12), row_y + Pt(2),
                       Inches(0.28), Inches(0.3), color)
            box(s, cx + Inches(0.12), row_y + Pt(2), Inches(0.28), Inches(0.3),
                text=lbl, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            # Text
            box(s, cx + Inches(0.48), row_y, Inches(5.5), Inches(0.48),
                text=data, font_size=10, color=result_colors[i], wrap=True)

    return s


# ---------------------------------------------------------------------------
# Slide 5 — Gains, Losses & Reflection
# ---------------------------------------------------------------------------
def slide_05_reflection(prs):
    s = blank_slide(prs)
    bg(s)

    box(s, Inches(0.5), Inches(0.22), Inches(6), Inches(0.38),
        text='04  /  Gains, Losses & Reflection', font_size=11, color=GOLD, bold=True)
    box(s, Inches(0.5), Inches(0.6), Inches(12), Inches(0.72),
        text='Looking Back — Honest Self-Assessment', font_size=30, bold=True, color=WHITE)
    accent_line(s, Inches(1.42))

    # Left — Gains
    shape_rect(s, Inches(0.4), Inches(1.6), Inches(6.0), Inches(5.65), DARK_CARD,
               line_color=GREEN, line_width=1.5)
    shape_rect(s, Inches(0.4), Inches(1.6), Inches(6.0), Inches(0.42), GREEN)
    box(s, Inches(0.55), Inches(1.65), Inches(5.7), Inches(0.35),
        text='✓  What I Gained', font_size=14, bold=True, color=WHITE)

    gains = [
        ('Built team-scale solutions',
         'Focused on tools that support the entire HR team, not just individual tasks.'),
        ('Deepened industry expertise',
         "Gained strong understanding of EdTech HR operations and 51Talk's business model."),
        ('Started learning Chinese',
         'Investing in language skills to connect with the broader 51Talk organization.'),
        ('Advanced AI & automation skills',
         'Leveraged AI tools to modernize HR workflows and build scalable systems.'),
    ]
    gy = Inches(2.18)
    for title, desc in gains:
        box(s, Inches(0.6), gy, Inches(5.6), Inches(0.32),
            text=f'▸  {title}', font_size=12, bold=True, color=GREEN)
        box(s, Inches(0.6), gy + Inches(0.32), Inches(5.6), Inches(0.42),
            text=f'   {desc}', font_size=11, color=LIGHT_GRAY, wrap=True)
        gy += Inches(0.92)

    # Right — What I'd Do Differently
    shape_rect(s, Inches(6.85), Inches(1.6), Inches(6.1), Inches(5.65), DARK_CARD,
               line_color=AMBER, line_width=1.5)
    shape_rect(s, Inches(6.85), Inches(1.6), Inches(6.1), Inches(0.42), AMBER)
    box(s, Inches(7.0), Inches(1.65), Inches(5.8), Inches(0.35),
        text="⟳  What I'd Do Differently", font_size=14, bold=True, color=NAVY)

    improvements = [
        ('Accelerate knowledge acquisition',
         "I'd use a more structured approach to ramp up faster on systems, "
         'processes, and people.'),
        ('Build management trust earlier',
         "I'd prioritize more frequent check-ins and proactive communication "
         'from day one.'),
        ('Implement systems sooner',
         'The automation and iTalent projects showed clear value — '
         'earlier implementation means earlier impact.'),
        ('Invest more in Chinese',
         'Language learning is a priority I want to pursue more consistently '
         'going forward.'),
    ]
    iy = Inches(2.18)
    for title, desc in improvements:
        box(s, Inches(7.0), iy, Inches(5.8), Inches(0.32),
            text=f'▸  {title}', font_size=12, bold=True, color=AMBER)
        box(s, Inches(7.0), iy + Inches(0.32), Inches(5.8), Inches(0.48),
            text=f'   {desc}', font_size=11, color=LIGHT_GRAY, wrap=True)
        iy += Inches(0.95)

    return s


# ---------------------------------------------------------------------------
# Slide 6 — Next Stage Plan & Support Needed
# ---------------------------------------------------------------------------
def slide_06_next(prs):
    s = blank_slide(prs)
    bg(s)

    box(s, Inches(0.5), Inches(0.22), Inches(6), Inches(0.38),
        text='05  /  Next Stage Plan', font_size=11, color=GOLD, bold=True)
    box(s, Inches(0.5), Inches(0.6), Inches(12), Inches(0.72),
        text='Looking Forward — Priorities & Support Needed', font_size=28, bold=True, color=WHITE)
    accent_line(s, Inches(1.42))

    box(s, Inches(0.5), Inches(1.6), Inches(5), Inches(0.38),
        text='Next Stage Priorities', font_size=14, bold=True, color=GOLD)

    priorities = [
        ('Complete iTalent ESS Portal Rollout',
         'Finalize configuration, run CC Team pilot, and transition the full company '
         'to digital leave management.'),
        ('HR Analytics',
         'Build comprehensive HR analytics reporting covering turnover, headcount, '
         'attendance, and performance trends.'),
        ('AI & Data Projects for HR',
         'Expand automation — smarter leave processing, predictive analytics, '
         'automated reporting, and intelligent HR workflows.'),
        ('Continue Learning Chinese',
         'Invest consistently in language learning to strengthen collaboration '
         'with the broader 51Talk organization.'),
    ]

    py = Inches(2.1)
    for i, (title, desc) in enumerate(priorities):
        # Number badge
        shape_rect(s, Inches(0.5), py + Pt(2), Inches(0.35), Inches(0.35), GOLD)
        box(s, Inches(0.5), py + Pt(2), Inches(0.35), Inches(0.35),
            text=str(i + 1), font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        box(s, Inches(1.0), py, Inches(7.0), Inches(0.32),
            text=title, font_size=13, bold=True, color=WHITE)
        box(s, Inches(1.0), py + Inches(0.33), Inches(7.0), Inches(0.42),
            text=desc, font_size=11, color=LIGHT_GRAY, wrap=True)
        py += Inches(0.9)

    # Support needed card
    shape_rect(s, Inches(8.55), Inches(1.55), Inches(4.45), Inches(5.7), DARK_CARD,
               line_color=GOLD, line_width=1.5)
    shape_rect(s, Inches(8.55), Inches(1.55), Inches(4.45), Inches(0.42), GOLD)
    box(s, Inches(8.7), Inches(1.6), Inches(4.1), Inches(0.35),
        text='Support Needed', font_size=14, bold=True, color=NAVY)

    supports = [
        ('Claude AI — Max Subscription',
         'Required to continue building and scaling HR automation and AI projects '
         'that benefit the entire HR department.'),
        ('Broader System Access',
         'Access to relevant HR and operational systems to enable deeper analytics '
         'and automation capabilities.'),
    ]
    sy = Inches(2.15)
    for title, desc in supports:
        box(s, Inches(8.7), sy, Inches(4.1), Inches(0.35),
            text=f'▸  {title}', font_size=12, bold=True, color=GOLD)
        box(s, Inches(8.7), sy + Inches(0.36), Inches(4.1), Inches(0.85),
            text=desc, font_size=11, color=LIGHT_GRAY, wrap=True)
        sy += Inches(1.45)

    return s


# ---------------------------------------------------------------------------
# Slide 7 — Thank You
# ---------------------------------------------------------------------------
def slide_07_thanks(prs):
    s = blank_slide(prs)
    bg(s)

    # Decorative stars
    for (lx, ly, sz) in [
        (Inches(10.5), Inches(0.5), Inches(0.4)),
        (Inches(12.2), Inches(1.5), Inches(0.25)),
        (Inches(11.2), Inches(2.8), Inches(0.18)),
        (Inches(0.4),  Inches(1.0), Inches(0.22)),
        (Inches(1.8),  Inches(6.5), Inches(0.18)),
        (Inches(12.7), Inches(5.8), Inches(0.28)),
    ]:
        star_deco(s, lx, ly, sz)

    # Left gold bar
    shape_rect(s, Inches(0), Inches(0), Inches(0.1), H, GOLD)

    # Large decorative circle
    circ = s.shapes.add_shape(MSO_OVAL, Inches(9.8), Inches(0.8), Inches(4.2), Inches(4.2))
    circ.fill.solid(); circ.fill.fore_color.rgb = DARK_CARD
    circ.line.color.rgb = GOLD; circ.line.width = Pt(2)

    box(s, Inches(0.5), Inches(2.0), Inches(9), Inches(1.5),
        text='Thank You', font_size=64, bold=True, color=WHITE)

    accent_line(s, Inches(3.75), width=Inches(5.5))

    box(s, Inches(0.5), Inches(4.05), Inches(9), Inches(0.55),
        text='Ahmed Al-Sadeq  |  HR Business Partner  |  51Talk Egypt',
        font_size=16, color=LIGHT_GRAY)

    box(s, Inches(0.5), Inches(4.65), Inches(8), Inches(0.5),
        text='Open for questions & discussion',
        font_size=15, color=GOLD)

    # Brand label bottom right
    box(s, Inches(10.5), Inches(6.9), Inches(2.5), Inches(0.38),
        text='51Talk Egypt', font_size=11, color=GOLD,
        align=PP_ALIGN.RIGHT)

    return s


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = new_prs()
    slide_01_cover(prs)
    slide_02_intro(prs)
    slide_03_performance(prs)
    slide_04_values(prs)
    slide_05_reflection(prs)
    slide_06_next(prs)
    slide_07_thanks(prs)

    out = 'docs/Probation-Review-Ahmed-AlSadeq.pptx'
    prs.save(out)
    print(f'Saved: {out}')


if __name__ == '__main__':
    main()
