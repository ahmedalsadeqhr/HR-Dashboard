"""
Convert iTalent planning documents to Word (.docx) and PowerPoint (.pptx)
"""
from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PPTXColor
from pptx.enum.text import PP_ALIGN
import re

# ─── Colors ───────────────────────────────────────────────────────────────────
BRAND_BLUE   = RGBColor(0x00, 0x56, 0xA8)
BRAND_GRAY   = RGBColor(0x5A, 0x5A, 0x5A)
TABLE_HEADER = RGBColor(0x00, 0x56, 0xA8)
TABLE_ALT    = RGBColor(0xF0, 0xF6, 0xFF)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)

PPTX_BLUE  = PPTXColor(0x00, 0x56, 0xA8)
PPTX_GRAY  = PPTXColor(0x5A, 0x5A, 0x5A)
PPTX_WHITE = PPTXColor(0xFF, 0xFF, 0xFF)
PPTX_ALT   = PPTXColor(0xF0, 0xF6, 0xFF)

# ─── Word helpers ──────────────────────────────────────────────────────────────
def set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)

def add_table(doc, headers, rows, col_widths=None):
    table = doc.add_table(rows=1+len(rows), cols=len(headers))
    table.style = 'Table Grid'
    # Header row
    hdr = table.rows[0]
    for i, h in enumerate(headers):
        cell = hdr.cells[i]
        cell.text = h
        set_cell_bg(cell, '0056A8')
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.runs[0]
        run.bold = True
        run.font.color.rgb = WHITE
        run.font.size = Pt(9)
    # Data rows
    for ri, row in enumerate(rows):
        tr = table.rows[ri+1]
        bg = 'F0F6FF' if ri % 2 == 0 else 'FFFFFF'
        for ci, val in enumerate(row):
            cell = tr.cells[ci]
            cell.text = str(val)
            set_cell_bg(cell, bg)
            cell.paragraphs[0].runs[0].font.size = Pt(9)
    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            for row in table.rows:
                row.cells[i].width = Inches(w)
    doc.add_paragraph()
    return table

def add_heading(doc, text, level=1, color=None):
    p = doc.add_heading(text, level=level)
    if color:
        for run in p.runs:
            run.font.color.rgb = color
    return p

def add_para(doc, text, bold=False, italic=False, size=10):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = bold
    run.italic = italic
    run.font.size = Pt(size)
    return p

# ─── Document 1: Email to Manager ─────────────────────────────────────────────
def build_email_manager():
    doc = Document()
    # Page margins
    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    # Header block
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run('51Talk Egypt — HR Department')
    run.bold = True
    run.font.size = Pt(13)
    run.font.color.rgb = BRAND_BLUE

    doc.add_paragraph()

    meta = [
        ('To:', '[Manager Name]'),
        ('From:', 'Ahmed Elsadek — HRBP, 51Talk Egypt'),
        ('Date:', 'March 1, 2026'),
        ('Subject:', 'iTalent HR System — Full Implementation Plan for Your Review'),
    ]
    for label, val in meta:
        p = doc.add_paragraph()
        r = p.add_run(label + '  ')
        r.bold = True
        r.font.color.rgb = BRAND_BLUE
        p.add_run(val)

    doc.add_paragraph()

    add_heading(doc, 'Dear [Manager Name],', level=2, color=BRAND_GRAY)
    doc.add_paragraph(
        'Following our recent discussion and your approval to proceed, I have finalized the complete '
        'implementation plan for transitioning our leave management and employee self-service processes '
        'to iTalent. Please find the full details below.'
    )

    # Section 1 — Background
    add_heading(doc, '1. Background', level=2, color=BRAND_BLUE)
    doc.add_paragraph(
        'The company has introduced iTalent as our new HR system. I have conducted a hands-on exploration '
        'of the live system and mapped its capabilities against our Employee Handbook and existing Attendance '
        'Dashboard. The plan below details exactly how we will configure the system, transition employees, and go live.'
    )

    # Section 2 — Current vs Target
    add_heading(doc, '2. Current Situation vs. Target State', level=2, color=BRAND_BLUE)
    add_table(doc,
        ['Area', 'Current', 'After iTalent'],
        [
            ('Leave requests', 'Email to manager + hr.egy@51talk.com', 'Digital submission via iTalent ESS portal'),
            ('Approval tracking', 'No system — email threads', 'Full audit trail in iTalent'),
            ('Leave balance visibility', 'HR-managed only', 'Employees see real-time balance'),
            ('Leave data for payroll', 'Manual leave sheet (Excel)', 'Auto-exported from iTalent on 21st monthly'),
            ('Approval enforcement', 'Honor-based', 'Configurable multi-level workflows'),
        ],
        col_widths=[1.8, 2.5, 2.5]
    )

    # Section 3 — Already in system
    add_heading(doc, '3. What Is Already in the System', level=2, color=BRAND_BLUE)
    doc.add_paragraph('After a live review of iTalent (51talk.italent.cn):')
    p = doc.add_paragraph()
    p.add_run('✅ Already configured:').bold = True
    for item in ['7 leave types (Sick, Annual, Unpaid, Maternity, Paternity, Bereavement, Miscarriage)',
                 'Employee Self-Service portal with My Approval, My Attendance, My Profile',
                 'Batch approval and Process Delegation for managers']:
        doc.add_paragraph(item, style='List Bullet')
    p = doc.add_paragraph()
    p.add_run('⚠️ Requires configuration before go-live:').bold = True
    for item in ['2 leave types unnamed (will be set as Casual Leave & Marriage Leave)',
                 'Annual Leave entitlement showing 0 days — needs rule setup',
                 'Approval workflows not yet assigned per leave type',
                 'Email notifications not configured',
                 'Employee data not yet imported']:
        doc.add_paragraph(item, style='List Bullet')

    # Section 4 — Confirmed Decisions
    add_heading(doc, '4. Your Confirmed Decisions (Already Incorporated)', level=2, color=BRAND_BLUE)
    add_table(doc,
        ['Decision', 'Your Answer'],
        [
            ('Pilot department', 'CC Team'),
            ('Approval chain setup', 'Add manually later'),
            ('Overtime', 'No overtime policy — feature will be disabled'),
            ('Annual Leave carryover', 'Allow carryover'),
            ('Issue Certificate types', 'HR Letter only'),
        ],
        col_widths=[2.5, 4.2]
    )

    # Section 5 — Plan
    add_heading(doc, '5. Implementation Plan — 7 Phases', level=2, color=BRAND_BLUE)
    add_table(doc,
        ['Phase', 'What', 'Duration'],
        [
            ('0 — Prerequisites', 'Admin access, employee data backup, CC Team manager briefed', '2–3 days'),
            ('1 — Leave Configuration', 'Fix leave types, set entitlements per handbook, configure carryover', 'Week 1'),
            ('2 — Approval Workflows', 'Build approval chains per leave type, set notification rules', 'Week 1'),
            ('3 — Dashboard Integration', 'Connect iTalent leave export to Attendance Dashboard', 'Week 2'),
            ('4 — ESS Configuration', 'HR Letter certificate, disable overtime, configure business trips', 'Week 2'),
            ('5 — Notifications', 'Email alerts for managers, employees, and HR on all leave events', 'Week 2'),
            ('6 — Pilot + Full Rollout', 'CC Team pilot → evaluate → company-wide rollout', 'Week 3–6'),
            ('7 — Operations', 'Monthly export-import cycle, balance management', 'Month 2+'),
        ],
        col_widths=[1.8, 4.0, 1.0]
    )

    # Section 6 — Rollout Strategy
    add_heading(doc, '6. Rollout Strategy', level=2, color=BRAND_BLUE)
    for line in ['Week 1–2:  Configure & integrate',
                 'Week 3:    CC Team pilot goes live (parallel with email process)',
                 'Week 4:    Review pilot, fix any gaps',
                 'Week 5–6:  Full company rollout + training sessions',
                 'End Week 6: Email-based leave process retired',
                 'Month 2+:  Full operations on iTalent']:
        doc.add_paragraph(line, style='List Bullet')
    doc.add_paragraph()
    p = doc.add_paragraph()
    p.add_run('No risk to employees during transition: ').bold = True
    p.add_run('Email-based process stays active during the pilot. No leave request will be lost. '
              'Opening leave balances will be imported before anyone accesses their account.')

    # Section 7 — Leave Entitlements
    add_heading(doc, '7. Leave Types & Entitlements (Per Handbook)', level=2, color=BRAND_BLUE)
    add_table(doc,
        ['Leave Type', 'Days', 'Pay', 'Key Rule'],
        [
            ('Annual Leave', '15 (<1yr) / 21 (≥1yr)', 'Full', 'Carryover allowed'),
            ('Casual Leave', '7/year, max 2/request', 'Full', 'Counts toward annual balance'),
            ('Sick Leave', 'As needed', '75%', "Doctor's note required"),
            ('Maternity Leave', '120 days', 'Full', 'After 1yr service, birth cert required'),
            ('Paternity Leave', '1 day × max 3/year', 'Full', 'Proof required'),
            ('Marriage Leave', '3 days', 'Full', 'Proof required'),
            ('Bereavement Leave', '3 days', 'Full', 'Parents/spouse/children only'),
            ('Miscarriage Leave', 'Per labor law', 'Full', 'Medical documentation required'),
            ('Unpaid Leave', 'As approved', '0%', 'Manager + HR dual approval'),
        ],
        col_widths=[1.5, 1.5, 0.6, 3.1]
    )

    # Section 8 — Risks
    add_heading(doc, '8. Risks & How We\'re Managing Them', level=2, color=BRAND_BLUE)
    add_table(doc,
        ['Risk', 'Mitigation'],
        [
            ('Employees don\'t adopt the system', 'Training sessions + Arabic/English quick-start guide'),
            ('Data mismatch with Attendance Dashboard', '2-week parallel run before cutover'),
            ('Manager approval delays', 'Process Delegation configured for all managers'),
            ('Leave history lost', 'Opening balances imported before any employee accesses the system'),
            ('System downtime', 'Backup: revert to email process temporarily'),
        ],
        col_widths=[2.8, 3.9]
    )

    # Section 9 — What I need
    add_heading(doc, '9. What I Still Need From You', level=2, color=BRAND_BLUE)
    for item in ['Pilot kickoff date — When should CC Team go live on iTalent?',
                 'IT system admin contact — I need to submit an admin access request to start configuration']:
        doc.add_paragraph('☐  ' + item)

    # Section 10 — Attachments
    add_heading(doc, '10. Attached Documents', level=2, color=BRAND_BLUE)
    add_table(doc,
        ['Document', 'Purpose'],
        [
            ('iTalent-Execution-Rollout-Plan', 'Full step-by-step transition plan with checklists'),
            ('iTalent-Implementation-Plan', 'Technical configuration details per phase'),
        ],
        col_widths=[2.5, 4.2]
    )

    # Closing
    doc.add_paragraph()
    doc.add_paragraph('Please let me know if you have any questions or would like to discuss any part of the plan.')
    doc.add_paragraph()
    doc.add_paragraph('Best regards,')
    p = doc.add_paragraph()
    p.add_run('Ahmed Elsadek').bold = True
    doc.add_paragraph('HRBP — 51Talk Egypt')
    doc.add_paragraph('hr.egy@51talk.com')

    path = r'C:\Users\high tech\Desktop\HRBP\docs\plans\Email-To-Manager-Full-Plan.docx'
    doc.save(path)
    print(f'Saved: {path}')


# ─── Document 2: Email to SysAdmin ────────────────────────────────────────────
def build_email_sysadmin():
    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run('51Talk Egypt — HR Department')
    run.bold = True
    run.font.size = Pt(13)
    run.font.color.rgb = BRAND_BLUE
    doc.add_paragraph()

    for label, val in [
        ('To:', '[System Admin Name / IT Team]'),
        ('From:', 'Ahmed Elsadek — HRBP, 51Talk Egypt'),
        ('Date:', 'March 1, 2026'),
        ('Subject:', 'iTalent Admin Access Request — HR Configuration Project'),
    ]:
        p = doc.add_paragraph()
        r = p.add_run(label + '  ')
        r.bold = True
        r.font.color.rgb = BRAND_BLUE
        p.add_run(val)

    doc.add_paragraph()
    add_heading(doc, 'Dear [System Admin Name],', level=2, color=BRAND_GRAY)
    doc.add_paragraph(
        'I am writing to request admin-level access to the iTalent HR system for 51Talk Egypt. '
        'We are initiating the configuration and rollout of iTalent for our leave management and '
        'employee self-service processes, and I need admin access to proceed with the setup.'
    )

    add_heading(doc, '1. What I Need Access To', level=2, color=BRAND_BLUE)
    doc.add_paragraph('Please grant admin/configuration access to the following modules in iTalent (51talk.italent.cn):')
    add_table(doc,
        ['Module', 'Access Required', 'Purpose'],
        [
            ('Leave Configuration', 'Read + Write', 'Rename leave types, set entitlements, configure carryover rules'),
            ('Workflow Configuration', 'Read + Write', 'Create and assign approval workflows per leave type'),
            ('Employee Management', 'Read + Write', 'Import employee master data and set opening leave balances'),
            ('Notification Settings', 'Read + Write', 'Configure email alerts for leave submissions, approvals, rejections'),
            ('Menu / UI Configuration', 'Read + Write', 'Fix typos in menu item names (2 locations)'),
            ('ESS Portal Settings', 'Read + Write', 'Enable/disable features (e.g., disable Overtime Application)'),
            ('Reports / Export', 'Read', 'Configure and test leave data export for Attendance Dashboard integration'),
        ],
        col_widths=[1.8, 1.3, 3.6]
    )

    add_heading(doc, '2. Access Account Details', level=2, color=BRAND_BLUE)
    add_table(doc,
        ['Field', 'Details'],
        [
            ('Requester Name', 'Ahmed Elsadek'),
            ('Role', 'HRBP'),
            ('Email', 'hr.egy@51talk.com'),
            ('Current iTalent Login', 'ahmedelsadek@51talk.com (or current employee login)'),
            ('Access Level Needed', 'HR Administrator / System Configurator'),
            ('Company/Tenant', '51Talk Egypt (51talk.italent.cn)'),
        ],
        col_widths=[2.2, 4.5]
    )

    add_heading(doc, '3. Why This Is Needed', level=2, color=BRAND_BLUE)
    doc.add_paragraph(
        'We have received management approval to configure iTalent for our leave management and employee '
        'self-service rollout. The project starts with a pilot on the CC Team and will expand company-wide '
        'over 6 weeks. Without admin access, I cannot:'
    )
    for item in [
        'Rename the 2 unnamed leave types ("Leave Item 6" and "Leave Item 34")',
        'Set leave entitlement rules (Annual Leave currently shows 0 days)',
        'Create approval workflows',
        'Import employee master data',
        'Configure notifications',
    ]:
        doc.add_paragraph(item, style='List Bullet')

    add_heading(doc, '4. Timeline & Urgency', level=2, color=BRAND_BLUE)
    add_table(doc,
        ['Milestone', 'Target Date'],
        [
            ('Admin access granted', 'ASAP — blocks all configuration work'),
            ('Stage 0 (prerequisites) complete', 'Within 2–3 days of access'),
            ('Stage 1 (configuration) start', 'Week 1 after access'),
            ('CC Team pilot go-live', 'Week 3'),
            ('Full company rollout', 'Week 5–6'),
        ],
        col_widths=[2.5, 4.2]
    )
    p = doc.add_paragraph()
    p.add_run('This access is on the critical path. ').bold = True
    p.add_run('Every day of delay pushes the pilot and company-wide rollout back by one day.')

    add_heading(doc, '5. Specific Configuration Items (for Your Reference)', level=2, color=BRAND_BLUE)
    doc.add_paragraph('Here is a summary of what I will configure once I have access — to help you assess the scope:')

    p = doc.add_paragraph()
    p.add_run('Leave Types (rename & configure):').bold = True
    for item in [
        '"Leave Item 6" → Casual Leave (7 days/year, max 2/request)',
        '"Leave Item 34" → Marriage Leave (3 days, proof required)',
        'Annual Leave → Set entitlement (15 days <1yr / 21 days ≥1yr, carryover allowed)',
        'Sick Leave → Set 75% pay rate, require attachment',
        'Maternity Leave → 120 days, 1yr service required',
        'Paternity Leave → 3 occurrences/year, 1 day each',
    ]:
        doc.add_paragraph(item, style='List Bullet')

    p = doc.add_paragraph()
    p.add_run('Approval Workflows (create 3 templates):').bold = True
    for item in ['Single level: Manager only', 'Two levels: Manager → HR', 'Three levels: Manager → Director → HR']:
        doc.add_paragraph(item, style='List Bullet')

    p = doc.add_paragraph()
    p.add_run('ESS Portal:').bold = True
    for item in [
        'Disable Overtime Application (no overtime policy)',
        'Issue Certificate: HR Letter only',
        'Fix menu typo: "Leave Applictaion" → "Leave Application" (2 locations)',
    ]:
        doc.add_paragraph(item, style='List Bullet')

    p = doc.add_paragraph()
    p.add_run('Employee Data:').bold = True
    for item in ['Import master data (names, IDs, departments, join dates)', 'Set opening leave balances per employee']:
        doc.add_paragraph(item, style='List Bullet')

    add_heading(doc, '6. Questions for You', level=2, color=BRAND_BLUE)
    for i, q in enumerate([
        'What is the process to request admin access in iTalent — do you submit directly to the iTalent vendor (51Talk HQ) or can you grant it locally?',
        'Is there an existing HR admin account I can use, or does a new account need to be created?',
        'Are there any restrictions on what HR can configure vs. what requires vendor involvement?',
        'Can you advise on how to import employee data — is there a standard import template in iTalent?',
    ], 1):
        doc.add_paragraph(f'{i}. {q}')

    doc.add_paragraph()
    doc.add_paragraph('Please let me know if you need any additional information or approvals to process this request.')
    doc.add_paragraph()
    doc.add_paragraph('Best regards,')
    p = doc.add_paragraph()
    p.add_run('Ahmed Elsadek').bold = True
    doc.add_paragraph('HRBP — 51Talk Egypt')
    doc.add_paragraph('hr.egy@51talk.com')
    doc.add_paragraph()
    doc.add_paragraph('CC: [Manager Name]')

    path = r'C:\Users\high tech\Desktop\HRBP\docs\plans\Email-To-SysAdmin-Access-Request.docx'
    doc.save(path)
    print(f'Saved: {path}')


# ─── Document 3: Technical Implementation Plan (.docx) ────────────────────────
def build_implementation_plan():
    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run('iTalent HR System — Implementation Plan')
    r.bold = True; r.font.size = Pt(16); r.font.color.rgb = BRAND_BLUE
    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = p2.add_run('51Talk Egypt  |  Prepared by: Ahmed Elsadek, HRBP  |  February 25, 2026')
    r2.font.size = Pt(10); r2.font.color.rgb = BRAND_GRAY
    doc.add_paragraph()
    doc.add_paragraph(
        'Goal: Configure iTalent for 51Talk Egypt\'s leave management and employee self-service '
        'aligned with company policy, Attendance Dashboard, and existing workflows.\n\n'
        'Architecture: iTalent (51talk.italent.cn) serves as the system of record for leave requests. '
        'The Attendance Dashboard (Streamlit) continues processing fingerprint data and generating '
        'penalty reports. iTalent bridges the gap by digitizing leave workflows and feeding approved '
        'leaves into the Dashboard.'
    )

    # Phase 1
    add_heading(doc, 'Phase 1: Leave Types Configuration', level=1, color=BRAND_BLUE)
    add_heading(doc, 'Task 1.1 — Audit & Rename Leave Types', level=2, color=BRAND_GRAY)
    doc.add_paragraph('Current State in iTalent:')
    add_table(doc,
        ['iTalent Name', 'Status'],
        [
            ('Sick Leave', '✅ Named correctly'),
            ('Unpaid Leave', '✅ Named correctly'),
            ('Annual Leave', '⚠️ Shows 0 days — entitlement not set'),
            ('Miscarriage Leave', '✅ Named correctly'),
            ('Bereavement Leave', '✅ Named correctly'),
            ('Maternity Leave', '✅ Named correctly'),
            ('Paternity Leave', '✅ Named correctly'),
            ('Leave Item 6', '❌ Unnamed — rename to Casual Leave'),
            ('Leave Item 34', '❌ Unnamed — rename to Marriage Leave'),
        ],
        col_widths=[2.5, 4.2]
    )
    doc.add_paragraph('Action Steps:')
    for item in [
        'Go to Admin → Attendance → Leave Configuration',
        'Rename "Leave Item 6" → Casual Leave',
        'Rename "Leave Item 34" → Marriage Leave',
        'Fix typo: "Leave Applictaion" → "Leave Application" (appears in 2 menu locations)',
    ]:
        doc.add_paragraph(item, style='List Number')

    add_heading(doc, 'Task 1.2 — Set Leave Entitlements Per Policy', level=2, color=BRAND_GRAY)
    add_table(doc,
        ['Leave Type', 'Entitlement', 'Pay', 'Conditions'],
        [
            ('Annual Leave', '15 days (<1yr) / 21 days (≥1yr)', 'Full', 'Carryover allowed'),
            ('Casual Leave', '7 days/year', 'Full', 'Max 2 days per request; counts toward annual leave'),
            ('Sick Leave', 'As needed', '75% salary', 'Requires stamped sick note from accredited hospital'),
            ('Maternity Leave', '120 days', 'Full', 'After 1 full year of service; requires birth certificate'),
            ('Paternity Leave', '1 day/occurrence (max 3×)', 'Full', 'Requires proof'),
            ('Marriage Leave', '3 days', 'Full', 'Requires proof; extra days from annual balance or unpaid'),
            ('Bereavement Leave', '3 days', 'Full', 'Parents, spouse, children only; proof required'),
            ('Miscarriage Leave', 'Per Egyptian Labor Law', 'Full', 'Medical documentation required'),
            ('Unpaid Leave', 'As approved', '0%', 'Manager + HR approval required'),
        ],
        col_widths=[1.5, 1.8, 0.7, 2.7]
    )
    doc.add_paragraph('Configuration Checklist:')
    for item in [
        'Annual Leave: Set balance rule (15 days <1yr / 21 days ≥1yr based on join date)',
        'Casual Leave: Set max 7 days/year, max 2 days per request',
        'Sick Leave: Set pay rate to 75%, require attachment (doctor\'s note)',
        'Maternity Leave: Set 120 days, require 1 year service eligibility',
        'Paternity Leave: Set max 3 occurrences/year, 1 day each',
        'Marriage Leave: Set 3 days, require attachment',
        'Bereavement Leave: Set 3 days, restricted beneficiaries',
        'Unpaid Leave: No balance limit, dual approval (manager + HR)',
    ]:
        doc.add_paragraph('☐  ' + item)

    # Phase 2
    add_heading(doc, 'Phase 2: Approval Workflow Configuration', level=1, color=BRAND_BLUE)
    add_heading(doc, 'Task 2.1 — Define Approval Chain Per Leave Type', level=2, color=BRAND_GRAY)
    doc.add_paragraph('Current process: Employee → Direct Manager → Higher Management → HR (hr.egy@51talk.com)')
    add_table(doc,
        ['Leave Type', 'Approval Steps'],
        [
            ('Annual Leave (≤ 3 days)', 'Step 1: Direct Manager'),
            ('Annual Leave (> 3 days)', 'Step 1: Direct Manager → Step 2: Director → Step 3: HR'),
            ('Casual Leave', 'Step 1: Direct Manager'),
            ('Sick Leave', 'Step 1: Direct Manager → Step 2: HR (documentation review)'),
            ('Maternity / Miscarriage', 'Step 1: Direct Manager → Step 2: HR'),
            ('Paternity', 'Step 1: Direct Manager'),
            ('Marriage Leave', 'Step 1: Direct Manager → Step 2: HR'),
            ('Bereavement Leave', 'Step 1: Direct Manager'),
            ('Unpaid Leave', 'Step 1: Direct Manager → Step 2: Director → Step 3: HR'),
        ],
        col_widths=[2.3, 4.4]
    )

    add_heading(doc, 'Task 2.2 — Configure Batch Approval Settings', level=2, color=BRAND_GRAY)
    for item in [
        'Enable batch approval for HR account (hr.egy@51talk.com)',
        'Enable Process Delegation for all manager-level roles',
        'Set notification: email to manager on new pending request',
    ]:
        doc.add_paragraph('☐  ' + item)

    # Phase 3
    add_heading(doc, 'Phase 3: Attendance Dashboard Integration', level=1, color=BRAND_BLUE)
    add_heading(doc, 'Task 3.1 — Align Leave Export Format', level=2, color=BRAND_GRAY)
    add_table(doc,
        ['Column', 'iTalent Field', 'Notes'],
        [
            ('CRM', 'Employee CRM ID', 'Must match master data'),
            ('Date', 'Leave Date', 'One row per day'),
            ('Leave Type', 'Leave Project name', 'e.g., Annual Leave, Sick Leave'),
            ('Status', 'Approved only', 'Only export approved leaves'),
        ],
        col_widths=[1.2, 2.0, 3.5]
    )
    doc.add_paragraph('Export Schedule: 21st of each month (aligns with Dashboard payroll period: 21st–20th). Format: Excel (.xlsx).')

    add_heading(doc, 'Task 3.2 — Leave Code Mapping', level=2, color=BRAND_GRAY)
    add_table(doc,
        ['iTalent Leave Type', 'Dashboard Status Code'],
        [
            ('Annual Leave', 'AL'), ('Casual Leave', 'CL'), ('Sick Leave', 'SL'),
            ('Unpaid Leave', 'UL'), ('Maternity Leave', 'ML'), ('Paternity Leave', 'PL'),
            ('Marriage Leave', 'MAL'), ('Bereavement Leave', 'BL'), ('Miscarriage Leave', 'MCL'),
        ],
        col_widths=[3.0, 2.0]
    )

    # Phase 4
    add_heading(doc, 'Phase 4: Employee Self-Service Configuration', level=1, color=BRAND_BLUE)
    add_table(doc,
        ['ESS Section', 'Sub-items', 'Action'],
        [
            ('My Approval', 'Pending, Processed, Applied', '✅ Already configured'),
            ('My Attendance', 'Leave Application, Overtime Application, My Business Trips', 'Disable Overtime'),
            ('Self-Service Application', 'Issue Certificate, Personnel Application', 'Configure: HR Letter only'),
            ('Common Functions', 'Leave Application shortcuts', '✅ Already visible'),
        ],
        col_widths=[1.8, 2.5, 2.4]
    )

    # Phase 5
    add_heading(doc, 'Phase 5: Notification & Communication Setup', level=1, color=BRAND_BLUE)
    add_table(doc,
        ['Event', 'Recipient', 'Email'],
        [
            ('New leave request submitted', 'Direct Manager', "Manager's email"),
            ('Leave approved', 'Employee + HR', 'hr.egy@51talk.com'),
            ('Leave rejected', 'Employee', 'Employee email'),
            ('Pending approval reminder (24hrs)', 'Manager', "Manager's email"),
            ('Leave request needs documentation', 'Employee', 'Employee email'),
        ],
        col_widths=[2.5, 1.8, 2.4]
    )

    # Phase 6
    add_heading(doc, 'Phase 6: Data Migration & Go-Live', level=1, color=BRAND_BLUE)
    add_heading(doc, 'Task 6.1 — Employee Master Data Import', level=2, color=BRAND_GRAY)
    add_table(doc,
        ['Field', 'Source', 'iTalent Field'],
        [
            ('AC-No / PS ID', 'Attendance Dashboard master', 'Employee ID'),
            ('CRM', 'Dashboard master', 'CRM'),
            ('Name', 'Dashboard master', 'Name'),
            ('Department', 'Dashboard master', 'Department'),
            ('Join Date', 'HR records', 'Join Date (annual leave entitlement)'),
            ('Vendor', 'Dashboard master', 'Vendor (Just HR / Migrate)'),
            ('National ID', 'HR records', 'National ID'),
        ],
        col_widths=[1.5, 2.3, 2.9]
    )

    add_heading(doc, 'Task 6.2 — Opening Balances Formula', level=2, color=BRAND_GRAY)
    for formula in [
        'Annual Leave (< 1yr): Opening balance = 15 − days_used_YTD',
        'Annual Leave (≥ 1yr): Opening balance = 21 − days_used_YTD',
        'Casual Leave: Opening balance = 7 − days_used_YTD',
        'Paternity: Opening occurrences = 3 − occurrences_used_YTD',
    ]:
        doc.add_paragraph(formula, style='List Bullet')

    add_heading(doc, 'Task 6.3 — Go-Live Checklist', level=2, color=BRAND_GRAY)
    for item in [
        'All leave types named and configured',
        'Approval workflows assigned and tested',
        'All employees imported with correct balances',
        'Email notifications tested end-to-end',
        'Leave export format validated with Attendance Dashboard',
        'HR and manager training completed',
        '"Leave Applictaion" typo fixed in both menu locations',
        'Leave Item 6 and Leave Item 34 renamed',
        'Annual Leave entitlement rule set (15/21 days based on tenure)',
        'Announcement sent to all employees',
    ]:
        doc.add_paragraph('☐  ' + item)

    # Phase 7
    add_heading(doc, 'Phase 7: Ongoing Operations', level=1, color=BRAND_BLUE)
    doc.add_paragraph('Monthly Workflow:')
    for item in [
        '1st–20th: Employees submit leaves via iTalent ESS',
        '21st: HR exports approved leaves from iTalent',
        '21st: Import leave export into Attendance Dashboard',
        '21st–25th: Dashboard generates penalty report',
        '30th: Payroll processed (basic salary)',
    ]:
        doc.add_paragraph(item, style='List Bullet')

    # Issues found
    add_heading(doc, 'Issues Found During Exploration', level=1, color=BRAND_BLUE)
    add_table(doc,
        ['#', 'Issue', 'Location', 'Priority'],
        [
            ('1', '"Leave Applictaion" typo', 'My Attendance sidebar + Common Functions', 'High'),
            ('2', '"Leave Item 6" unnamed', 'Leave type dropdown', 'High'),
            ('3', '"Leave Item 34" unnamed', 'Leave type dropdown', 'High'),
            ('4', 'Annual Leave shows 0 days', 'Leave entitlement config', 'High'),
            ('5', 'No Casual Leave configured', 'Leave types', 'High'),
            ('6', 'No Marriage Leave configured', 'Leave types', 'High'),
            ('7', 'Overtime Application form — fields unknown', 'My Attendance', 'Medium'),
            ('8', 'Certificate types not configured', 'Self-Service Application', 'Medium'),
        ],
        col_widths=[0.3, 2.2, 2.5, 1.0]
    )

    # Decisions
    add_heading(doc, 'Decisions Log — All Resolved ✅', level=1, color=BRAND_BLUE)
    add_table(doc,
        ['#', 'Question', 'Answer', 'Action'],
        [
            ('1', 'Approval chain — who are the managers?', 'To be added manually later', 'Skip pre-assignment; HR adds approvers per team post-pilot'),
            ('2', 'Overtime policy?', 'No overtime', 'Disable Overtime Application in ESS'),
            ('3', 'Annual Leave carryover?', 'Allow carryover', 'Configure iTalent to carry over unused annual leave'),
            ('4', 'Certificate types?', 'HR Letter only', 'Configure one certificate type in ESS'),
            ('5', 'Pilot department?', 'CC Team', 'Stage 4 targets CC Team'),
            ('6', 'Casual Leave vs Annual Leave deduction', 'Casual counts toward annual per handbook', 'Configure Casual Leave to deduct from Annual Leave balance'),
        ],
        col_widths=[0.3, 1.8, 1.5, 3.1]
    )

    path = r'C:\Users\high tech\Desktop\HRBP\docs\plans\iTalent-Implementation-Plan.docx'
    doc.save(path)
    print(f'Saved: {path}')


# ─── Document 4: Execution Rollout Plan (.pptx) ───────────────────────────────
def build_rollout_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # completely blank

    def add_slide(title_text, content_fn, subtitle=None):
        slide = prs.slides.add_slide(blank)
        # Blue header bar
        bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.2))
        bar.fill.solid(); bar.fill.fore_color.rgb = PPTX_BLUE
        bar.line.fill.background()
        # Title
        tf = bar.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title_text
        run.font.bold = True
        run.font.size = Pt(22)
        run.font.color.rgb = PPTX_WHITE
        tf.margin_left = Inches(0.3)
        tf.margin_top  = Inches(0.25)
        if subtitle:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            r2 = p2.add_run()
            r2.text = subtitle
            r2.font.size = Pt(11)
            r2.font.color.rgb = PPTXColor(0xCC, 0xE0, 0xFF)
        content_fn(slide)
        return slide

    def pptx_table(slide, headers, rows, left, top, width, height, col_ratios=None):
        from pptx.util import Emu
        cols = len(headers)
        tbl = slide.shapes.add_table(len(rows)+1, cols, left, top, width, height).table
        if col_ratios:
            total = sum(col_ratios)
            for i, r in enumerate(col_ratios):
                tbl.columns[i].width = int(width * r / total)
        # Header
        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = h
            cell.fill.solid(); cell.fill.fore_color.rgb = PPTX_BLUE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.bold = True; run.font.color.rgb = PPTX_WHITE; run.font.size = Pt(10)
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri+1, ci)
                cell.text = str(val)
                if ri % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = PPTX_ALT
                else:
                    cell.fill.solid(); cell.fill.fore_color.rgb = PPTX_WHITE
                p = cell.text_frame.paragraphs[0]
                p.runs[0].font.size = Pt(9)
        return tbl

    def add_bullets(slide, items, left, top, width, height, size=13, bold_first=False):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame; tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = item
            run.font.size = Pt(size)
            if bold_first and i == 0:
                run.font.bold = True

    def label(slide, text, left, top, width=Inches(2), height=Inches(0.4), size=11, bold=False, color=None):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        if color: r.font.color.rgb = color

    # ── Slide 1: Title ──
    def s1_content(slide):
        tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = 'iTalent Execution & Rollout Plan'
        r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = PPTX_BLUE
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = '51Talk Egypt  ·  HRBP: Ahmed Elsadek  ·  March 2026'
        r2.font.size = Pt(14); r2.font.color.rgb = PPTX_GRAY
        # decorative line
        from pptx.util import Pt as PPt
        line = slide.shapes.add_shape(1, Inches(3), Inches(4.5), Inches(7), Inches(0.05))
        line.fill.solid(); line.fill.fore_color.rgb = PPTX_BLUE
        line.line.fill.background()
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.7), Inches(11), Inches(0.8))
        tf2 = tb2.text_frame
        p3 = tf2.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = 'Transitioning 51Talk Egypt leave management from email to iTalent ESS Portal'
        r3.font.size = Pt(13); r3.font.color.rgb = PPTX_GRAY

    add_slide('', s1_content)

    # ── Slide 2: Guiding Principles ──
    def s2_content(slide):
        principles = [
            ('🚫 No Big Bang', 'Pilot first, fix issues, then expand'),
            ('📂 No Data Loss', 'All historical leave balances imported before go-live'),
            ('⚡ No Disruption', 'Parallel run ensures no employee is left without a process'),
            ('🔔 No Surprises', 'Employees and managers trained before they see the system'),
        ]
        for i, (title, desc) in enumerate(principles):
            col = i % 2; row = i // 2
            left = Inches(0.5 + col * 6.4)
            top  = Inches(1.5 + row * 2.5)
            box = slide.shapes.add_shape(1, left, top, Inches(6.0), Inches(2.2))
            box.fill.solid(); box.fill.fore_color.rgb = PPTX_ALT
            box.line.color.rgb = PPTX_BLUE
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = title
            r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = PPTX_BLUE
            p2 = tf.add_paragraph()
            r2 = p2.add_run(); r2.text = desc
            r2.font.size = Pt(11); r2.font.color.rgb = PPTX_GRAY

    add_slide('Guiding Principles', s2_content, 'Four principles to ensure a safe transition')

    # ── Slide 3: Timeline ──
    def s3_content(slide):
        weeks = [
            ('Week 1', 'Configure & Setup', PPTX_BLUE),
            ('Week 2', 'Integrate & Notify', PPTXColor(0x00, 0x80, 0xC0)),
            ('Week 3', 'Pilot Go-Live', PPTXColor(0x00, 0xA0, 0x60)),
            ('Week 4', 'Evaluate & Fix', PPTXColor(0xE0, 0x80, 0x00)),
            ('Week 5–6', 'Full Rollout', PPTXColor(0xC0, 0x20, 0x20)),
            ('Month 2+', 'Operations', PPTXColor(0x60, 0x60, 0x60)),
        ]
        box_w = Inches(1.9)
        gap = Inches(0.12)
        start_x = Inches(0.4)
        for i, (label_text, desc, color) in enumerate(weeks):
            left = start_x + i * (box_w + gap)
            box = slide.shapes.add_shape(1, left, Inches(1.5), box_w, Inches(3.5))
            box.fill.solid(); box.fill.fore_color.rgb = color
            box.line.fill.background()
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.1); tf.margin_top = Inches(0.15)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = label_text
            r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = PPTX_WHITE
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run(); r2.text = desc
            r2.font.size = Pt(10); r2.font.color.rgb = PPTX_WHITE
        # Arrow
        arrow_top = Inches(3.0)
        a = slide.shapes.add_shape(1, Inches(0.3), arrow_top, Inches(12.7), Inches(0.08))
        a.fill.solid(); a.fill.fore_color.rgb = PPTX_GRAY
        a.line.fill.background()

    add_slide('6-Week Rollout Timeline', s3_content, 'Phased, low-risk approach from configuration to full operations')

    # ── Slide 4: Stage 0 — Prerequisites ──
    def s4_content(slide):
        pptx_table(slide,
            ['Task', 'Owner', 'Duration'],
            [
                ('Confirm iTalent admin access', 'HRBP + IT Admin', '2–3 days'),
                ('Export employee master data from Attendance Dashboard', 'HRBP', 'Day 1'),
                ('Collect org chart / reporting lines for all teams', 'HRBP', 'Day 1'),
                ('Pilot department confirmed: CC Team', 'Management ✅', 'Completed'),
                ('Get manager email list for approval chain setup', 'HRBP', 'Day 2'),
                ('Backup current leave sheet used in Attendance Dashboard', 'HRBP', 'Day 1'),
            ],
            Inches(0.4), Inches(1.4), Inches(12.5), Inches(3.0),
            col_ratios=[5, 2, 1.5]
        )
        add_bullets(slide, [
            '🚨  Admin access is on the critical path — every day of delay pushes pilot and rollout back by 1 day.'
        ], Inches(0.4), Inches(4.7), Inches(12.5), Inches(0.8), size=11)

    add_slide('Stage 0: Prerequisites  (2–3 Days)', s4_content, 'Must complete before any configuration begins')

    # ── Slide 5: Stage 1 — Leave Types ──
    def s5_content(slide):
        pptx_table(slide,
            ['Action', 'In iTalent', 'Expected Result'],
            [
                ('Rename "Leave Item 6"', 'Admin → Attendance → Leave Config', 'Name becomes "Casual Leave"'),
                ('Rename "Leave Item 34"', 'Admin → Attendance → Leave Config', 'Name becomes "Marriage Leave"'),
                ('Fix "Leave Applictaion" typo', 'Admin → Menu Config (2 locations)', 'Correct spelling in sidebar'),
                ('Set Annual Leave entitlement', 'Leave Config → Annual Leave → Balance Rule', '15 days (<1yr) / 21 days (≥1yr)'),
                ('Set Sick Leave pay rate', 'Leave Config → Sick Leave', '75% salary, attachment required'),
                ('Set Maternity Leave', 'Leave Config → Maternity', '120 days, 1yr service required'),
                ('Set Paternity Leave', 'Leave Config → Paternity', '3 occurrences max, 1 day each'),
                ('Set Unpaid Leave', 'Leave Config → Unpaid', 'No limit, dual approval'),
            ],
            Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.0),
            col_ratios=[3, 3.5, 3]
        )

    add_slide('Stage 1A: Fix Leave Types  (Day 1–2)', s5_content, 'Rename placeholders, set entitlements, fix UI typos')

    # ── Slide 6: Stage 1 — Workflows ──
    def s6_content(slide):
        workflows = [
            ('Workflow A\n"Single Approval"', 'Manager Only', 'Casual Leave\nPaternity\nBereavement\nBusiness Trips', PPTX_BLUE),
            ('Workflow B\n"HR Reviewed"', 'Manager → HR', 'Sick Leave\nMaternity\nMarriage\nMiscarriage', PPTXColor(0x00, 0x90, 0xC0)),
            ('Workflow C\n"Full Chain"', 'Manager → Director → HR', 'Annual Leave > 3 days\nUnpaid Leave', PPTXColor(0x00, 0x70, 0x50)),
        ]
        for i, (name, chain, types, color) in enumerate(workflows):
            left = Inches(0.4 + i * 4.3)
            box = slide.shapes.add_shape(1, left, Inches(1.5), Inches(4.1), Inches(4.0))
            box.fill.solid(); box.fill.fore_color.rgb = color
            box.line.fill.background()
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.15)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = name
            r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = PPTX_WHITE
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run(); r2.text = chain
            r2.font.size = Pt(11); r2.font.color.rgb = PPTXColor(0xCC, 0xE8, 0xFF)
            p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
            r3 = p3.add_run(); r3.text = '\nLeave Types:\n' + types
            r3.font.size = Pt(10); r3.font.color.rgb = PPTX_WHITE

    add_slide('Stage 1B: Approval Workflows  (Day 2–3)', s6_content, 'Three workflow templates covering all leave types')

    # ── Slide 7: Stage 2 — Integration ──
    def s7_content(slide):
        pptx_table(slide,
            ['Step', 'Action', 'Validation'],
            [
                ('1', 'Configure iTalent leave export (vertical format: CRM | Date | Leave Type | Status)', 'Export file generates correctly'),
                ('2', 'Export sample approved leaves from iTalent', 'Approved only, correct date range'),
                ('3', 'Import into Attendance Dashboard leave sheet', 'No import errors'),
                ('4', 'Run Dashboard report — verify penalties correct', '0 discrepancies vs manual sheet'),
                ('5', 'Confirm leave code mapping (AL, SL, CL, ML, PL, UL, BL, MAL, MCL)', 'All codes recognized'),
                ('6', 'Document Monthly SOP for HR', 'SOP reviewed and approved'),
            ],
            Inches(0.4), Inches(1.4), Inches(12.5), Inches(3.5),
            col_ratios=[0.5, 5, 3.5]
        )
        add_bullets(slide, [
            '📅  Export trigger: Manual — HR runs on 21st of each month  |  Format: Excel (.xlsx)  |  Date range: 21st prev → 20th current'
        ], Inches(0.4), Inches(5.2), Inches(12.5), Inches(0.8), size=10)

    add_slide('Stage 2: Attendance Dashboard Integration  (Week 2)', s7_content, 'Connect iTalent leave exports to the Streamlit Attendance Dashboard')

    # ── Slide 8: Stage 3 — Data Migration ──
    def s8_content(slide):
        pptx_table(slide,
            ['Field', 'Source', 'Notes'],
            [
                ('Employee ID (AC-No / PS ID)', 'Attendance Dashboard', 'Must match Dashboard master'),
                ('CRM', 'Attendance Dashboard', 'Links leave records to attendance'),
                ('Full Name', 'Attendance Dashboard', '—'),
                ('Department', 'Attendance Dashboard', '—'),
                ('Join Date', 'HR records', 'Critical — drives annual leave entitlement (15 vs 21 days)'),
                ('Vendor', 'Attendance Dashboard', 'Just HR / Migrate'),
                ('National ID', 'HR records', '—'),
                ('Manager / Approver', 'Org chart', 'Linked to workflow routing'),
                ('Email Address', 'Company directory', 'Needed for notification delivery'),
            ],
            Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.3),
            col_ratios=[2.5, 2.5, 4]
        )

    add_slide('Stage 3: Employee Data Migration  (Week 2)', s8_content, 'Import master data and set opening leave balances before any employee logs in')

    # ── Slide 9: Stage 4 — Pilot ──
    def s9_content(slide):
        days = [
            ('Day 1\nMonday', 'Kickoff\nBrief pilot manager\nShare employee quick-start guide\nEnable CC Team in iTalent'),
            ('Day 2–3', 'Live Usage\nEmployees submit pending leave requests\nManager approves via My Approval\nHRBP monitors for issues'),
            ('Day 4', 'Mid-Pilot Check\nReview all submissions\nVerify: notifications, approvals, balances\nFix any issues found'),
            ('Day 5\nFriday', 'Pilot Review\nDebrief with pilot manager\nDocument all issues\nGo / No-Go decision for full rollout'),
        ]
        for i, (day, content) in enumerate(days):
            left = Inches(0.4 + i * 3.2)
            box = slide.shapes.add_shape(1, left, Inches(1.5), Inches(3.0), Inches(4.5))
            box.fill.solid(); box.fill.fore_color.rgb = PPTX_ALT
            box.line.color.rgb = PPTX_BLUE
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.1); tf.margin_top = Inches(0.1)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = day
            r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = PPTX_BLUE
            for line in content.split('\n'):
                p2 = tf.add_paragraph()
                r2 = p2.add_run(); r2.text = '• ' + line if line else ''
                r2.font.size = Pt(10); r2.font.color.rgb = PPTX_GRAY

    add_slide('Stage 4: Pilot — CC Team  (Week 3)', s9_content, 'Validate full workflow in real conditions before company-wide rollout')

    # ── Slide 10: Stage 5 & 6 — Rollout ──
    def s10_content(slide):
        pptx_table(slide,
            ['Session', 'Audience', 'Duration', 'Content'],
            [
                ('Manager Training', 'All team managers', '1 hour', 'My Approval, Process Delegation, Batch Approval'),
                ('Employee Training', 'All staff', '30 min', 'Leave Application, My Attendance, My Profile'),
                ('HR Operations', 'HR team', '1 hour', 'Monthly export process, balance management, reports'),
            ],
            Inches(0.4), Inches(1.5), Inches(12.5), Inches(2.0),
            col_ratios=[2, 2, 1, 4]
        )
        add_bullets(slide, [
            'Week 5: Training for all managers (Day 1) → Training for all employees by department (Day 2–3) → Employees begin submitting in iTalent (Day 4–5)',
            'Week 6: Parallel run — employees submit in BOTH email + iTalent. HR approves in both. End of week: compare logs → confirm consistency.',
            'End of Week 6: Email process retired. HR email auto-reply updated. All pending requests migrated to iTalent.',
        ], Inches(0.4), Inches(3.8), Inches(12.5), Inches(2.5), size=11)

    add_slide('Stage 5 & 6: Fix, Train & Full Rollout  (Week 4–6)', s10_content, 'Training sessions, parallel run, and cutover from email to iTalent')

    # ── Slide 11: Risk Register ──
    def s11_content(slide):
        pptx_table(slide,
            ['Risk', 'Likelihood', 'Impact', 'Mitigation'],
            [
                ('Employees don\'t know how to use the system', 'High', 'Medium', 'Training + quick-start guide + HR support'),
                ('Manager approval bottleneck', 'Medium', 'High', 'Process Delegation configured; 48hr escalation'),
                ('Leave balance errors at import', 'Medium', 'High', 'Verify 5 sample employees before import'),
                ('Dashboard not recognizing iTalent leave codes', 'Low', 'High', 'Test export-import before pilot starts'),
                ('Employee not in iTalent system', 'Low', 'Medium', 'Full audit of employee list before go-live'),
                ('Manager rejects system adoption', 'Low', 'High', 'Escalate to Director; demo system value'),
                ('System downtime during critical period', 'Low', 'Medium', 'Backup: revert to email process temporarily'),
            ],
            Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.5),
            col_ratios=[4, 1.5, 1.2, 4]
        )

    add_slide('Risk Register', s11_content, 'Known risks and mitigations for a safe transition')

    # ── Slide 12: Success Metrics ──
    def s12_content(slide):
        pptx_table(slide,
            ['Metric', 'Target', 'Measurement'],
            [
                ('Leave requests via iTalent', '100% within 4 weeks of cutover', 'Count of email requests received after cutover'),
                ('Approval turnaround time', '< 24 hours', 'iTalent approval timestamps'),
                ('Leave data accuracy in Dashboard', '0 discrepancies', 'Compare iTalent export vs old leave sheet'),
                ('Employee satisfaction with ESS', '> 80% positive', 'Short survey after 1 month'),
                ('HR time saved on leave admin', '> 2 hours/week', 'HR self-reported'),
            ],
            Inches(0.4), Inches(1.4), Inches(12.5), Inches(3.5),
            col_ratios=[3.5, 2.5, 4]
        )
        add_bullets(slide, [
            'Contact: Ahmed Elsadek | hr.egy@51talk.com  ·  System: 51talk.italent.cn  ·  Based on: 51Talk Egypt Employee Handbook + Attendance Dashboard v2.2'
        ], Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.6), size=10)

    add_slide('Success Metrics & Contact', s12_content, 'How we measure a successful rollout')

    path = r'C:\Users\high tech\Desktop\HRBP\docs\plans\iTalent-Execution-Rollout-Plan.pptx'
    prs.save(path)
    print(f'Saved: {path}')


# ─── Run all ───────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    print('Generating documents...')
    build_email_manager()
    build_email_sysadmin()
    build_implementation_plan()
    build_rollout_pptx()
    print('All done!')
